from __future__ import annotations

import re
from dataclasses import dataclass
from datetime import datetime
from hashlib import sha1
from io import BytesIO
from pathlib import Path

import httpx

from app.models import DeckBlueprint, RetrievalResult


INTENT_COLORS: dict[str, tuple[int, int, int]] = {
    "technical": (29, 78, 216),
    "business": (22, 101, 52),
    "academic": (55, 65, 81),
    "creative": (194, 65, 12),
}


@dataclass
class RenderedPresentationFiles:
    deck_path: Path
    notes_path: Path
    background_image: str | None = None


class PptxDeckBuilder:
    def __init__(self, output_dir: Path, backgrounds_dir: Path | None = None) -> None:
        self.output_dir = output_dir
        self.backgrounds_dir = backgrounds_dir

    def build(
        self,
        *,
        deck: DeckBlueprint,
        topic: str,
        intent: str,
        context: list[RetrievalResult] | None = None,
    ) -> RenderedPresentationFiles:
        try:
            from pptx import Presentation
            from pptx.dml.color import RGBColor
            from pptx.enum.shapes import MSO_SHAPE
            from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
            from pptx.util import Inches, Pt
        except ImportError as exc:
            raise RuntimeError(
                "python-pptx is not installed. Run `uv pip install -r requirements.txt` in Backend."
            ) from exc

        self.output_dir.mkdir(parents=True, exist_ok=True)
        background_path = self._select_background_path(topic)
        accent = INTENT_COLORS.get(intent, (29, 78, 216))
        context_map = {self._orb_id(item): item for item in context or []}
        used_visual_urls: set[str] = set()
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        slug = self._slugify(topic)

        deck_prs = Presentation()
        deck_prs.slide_width = Inches(13.333)
        deck_prs.slide_height = Inches(7.5)

        notes_prs = Presentation()
        notes_prs.slide_width = Inches(13.333)
        notes_prs.slide_height = Inches(7.5)

        for index, slide_plan in enumerate(deck.slides):
            slide = deck_prs.slides.add_slide(deck_prs.slide_layouts[6])
            self._paint_background(
                slide=slide,
                accent=accent,
                slide_width=deck_prs.slide_width,
                slide_height=deck_prs.slide_height,
                background_path=background_path,
                soften_overlay=0.18,
            )
            self._render_main_slide(
                slide=slide,
                slide_plan=slide_plan,
                accent=accent,
                index=index,
                context_map=context_map,
                used_visual_urls=used_visual_urls,
            )

            notes_slide = notes_prs.slides.add_slide(notes_prs.slide_layouts[6])
            self._paint_background(
                slide=notes_slide,
                accent=accent,
                slide_width=notes_prs.slide_width,
                slide_height=notes_prs.slide_height,
                background_path=background_path,
                soften_overlay=0.1,
            )
            self._render_notes_slide(
                slide=notes_slide,
                slide_plan=slide_plan,
                accent=accent,
            )

        deck_path = self.output_dir / f"{slug}_deck_{timestamp}.pptx"
        notes_path = self.output_dir / f"{slug}_speaker_notes_{timestamp}.pptx"
        deck_prs.save(deck_path)
        notes_prs.save(notes_path)
        return RenderedPresentationFiles(
            deck_path=deck_path,
            notes_path=notes_path,
            background_image=background_path.name if background_path else None,
        )

    def _render_main_slide(
        self,
        *,
        slide,
        slide_plan,
        accent: tuple[int, int, int],
        index: int,
        context_map: dict[str, RetrievalResult],
        used_visual_urls: set[str],
    ) -> None:
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
        from pptx.util import Inches, Pt

        section_chip = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.72),
            Inches(0.42),
            Inches(1.55),
            Inches(0.38),
        )
        section_chip.fill.solid()
        section_chip.fill.fore_color.rgb = RGBColor(*accent)
        section_chip.fill.transparency = 0.08
        section_chip.line.fill.background()
        chip_frame = section_chip.text_frame
        chip_run = chip_frame.paragraphs[0].add_run()
        chip_run.text = f"SLIDE {slide_plan.slide_number:02d}"
        chip_run.font.size = Pt(10)
        chip_run.font.bold = True
        chip_run.font.color.rgb = RGBColor(255, 255, 255)

        content_panel = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.42),
            Inches(0.92),
            Inches(7.2),
            Inches(5.5),
        )
        content_panel.fill.solid()
        content_panel.fill.fore_color.rgb = RGBColor(255, 255, 255)
        content_panel.fill.transparency = 0.14
        content_panel.line.fill.background()

        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.1), Inches(6.4), Inches(0.9))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_run = title_frame.paragraphs[0].add_run()
        title_run.text = slide_plan.title
        title_run.font.size = Pt(24 if index else 28)
        title_run.font.bold = True
        title_run.font.color.rgb = RGBColor(15, 23, 42)

        objective_box = slide.shapes.add_textbox(Inches(0.76), Inches(1.72), Inches(6.3), Inches(0.48))
        objective_frame = objective_box.text_frame
        objective_run = objective_frame.paragraphs[0].add_run()
        objective_run.text = slide_plan.objective
        objective_run.font.size = Pt(12)
        objective_run.font.bold = True
        objective_run.font.color.rgb = RGBColor(*accent)

        summary_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(6.0), Inches(1.42))
        summary_frame = summary_box.text_frame
        summary_frame.word_wrap = True
        summary_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        summary_run = summary_frame.paragraphs[0].add_run()
        summary_run.text = slide_plan.summary_paragraph
        summary_run.font.size = Pt(14)
        summary_run.font.color.rgb = RGBColor(51, 65, 85)

        bullet_box = slide.shapes.add_textbox(Inches(0.86), Inches(3.75), Inches(5.95), Inches(2.15))
        bullet_frame = bullet_box.text_frame
        bullet_frame.word_wrap = True
        bullet_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        for point_index, point in enumerate(slide_plan.key_points):
            paragraph = bullet_frame.paragraphs[0] if point_index == 0 else bullet_frame.add_paragraph()
            paragraph.text = point
            paragraph.level = 0
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = RGBColor(31, 41, 55)
            paragraph.space_after = Pt(8)

        image_placed = self._place_visual_asset(
            slide=slide,
            slide_plan=slide_plan,
            context_map=context_map,
            accent=accent,
            used_visual_urls=used_visual_urls,
        )
        if not image_placed:
            self._render_visual_brief(slide=slide, slide_plan=slide_plan, accent=accent)

        evidence_panel = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(7.78),
            Inches(6.15),
            Inches(4.86),
            Inches(0.5),
        )
        evidence_panel.fill.solid()
        evidence_panel.fill.fore_color.rgb = RGBColor(255, 255, 255)
        evidence_panel.fill.transparency = 0.12
        evidence_panel.line.fill.background()
        evidence_frame = evidence_panel.text_frame
        evidence_run = evidence_frame.paragraphs[0].add_run()
        evidence_run.text = "Evidence: " + (", ".join(slide_plan.evidence_orbs[:3]) if slide_plan.evidence_orbs else "reasoned synthesis")
        evidence_run.font.size = Pt(10)
        evidence_run.font.color.rgb = RGBColor(31, 41, 55)

        footer_box = slide.shapes.add_textbox(Inches(0.7), Inches(6.74), Inches(12.0), Inches(0.34))
        footer_frame = footer_box.text_frame
        footer_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        footer_run = footer_frame.paragraphs[0].add_run()
        footer_run.text = (
            f"Slide {slide_plan.slide_number}  |  Evidence: "
            + (", ".join(slide_plan.evidence_orbs) if slide_plan.evidence_orbs else "reasoned synthesis")
        )
        footer_run.font.size = Pt(10)
        footer_run.font.color.rgb = RGBColor(71, 85, 105)

    def _render_notes_slide(self, *, slide, slide_plan, accent: tuple[int, int, int]) -> None:
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import MSO_AUTO_SIZE
        from pptx.util import Inches, Pt

        eyebrow = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.72),
            Inches(0.42),
            Inches(2.1),
            Inches(0.38),
        )
        eyebrow.fill.solid()
        eyebrow.fill.fore_color.rgb = RGBColor(*accent)
        eyebrow.fill.transparency = 0.08
        eyebrow.line.fill.background()
        eyebrow_frame = eyebrow.text_frame
        eyebrow_run = eyebrow_frame.paragraphs[0].add_run()
        eyebrow_run.text = "SPEAKER NOTES"
        eyebrow_run.font.size = Pt(10)
        eyebrow_run.font.bold = True
        eyebrow_run.font.color.rgb = RGBColor(255, 255, 255)

        title_box = slide.shapes.add_textbox(Inches(0.72), Inches(0.65), Inches(11.6), Inches(0.9))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_run = title_frame.paragraphs[0].add_run()
        title_run.text = f"Speaker Notes | Slide {slide_plan.slide_number}: {slide_plan.title}"
        title_run.font.size = Pt(24)
        title_run.font.bold = True
        title_run.font.color.rgb = RGBColor(255, 255, 255)

        note_panel = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.65),
            Inches(1.42),
            Inches(12.0),
            Inches(4.95),
        )
        note_panel.fill.solid()
        note_panel.fill.fore_color.rgb = RGBColor(255, 255, 255)
        note_panel.fill.transparency = 0.12
        note_panel.line.fill.background()

        summary_box = slide.shapes.add_textbox(Inches(0.95), Inches(1.82), Inches(10.8), Inches(1.2))
        summary_frame = summary_box.text_frame
        summary_frame.word_wrap = True
        summary_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        summary_title = summary_frame.paragraphs[0].add_run()
        summary_title.text = "Slide summary"
        summary_title.font.size = Pt(13)
        summary_title.font.bold = True
        summary_title.font.color.rgb = RGBColor(*accent)
        summary_body = summary_frame.add_paragraph()
        summary_body.text = slide_plan.summary_paragraph
        summary_body.font.size = Pt(16)
        summary_body.font.color.rgb = RGBColor(31, 41, 55)

        notes_box = slide.shapes.add_textbox(Inches(0.95), Inches(3.08), Inches(10.8), Inches(2.2))
        notes_frame = notes_box.text_frame
        notes_frame.word_wrap = True
        notes_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        notes_title = notes_frame.paragraphs[0].add_run()
        notes_title.text = "Speaker notes"
        notes_title.font.size = Pt(13)
        notes_title.font.bold = True
        notes_title.font.color.rgb = RGBColor(*accent)
        notes_body = notes_frame.add_paragraph()
        notes_body.text = slide_plan.speaker_notes
        notes_body.font.size = Pt(16)
        notes_body.font.color.rgb = RGBColor(31, 41, 55)

        orb_box = slide.shapes.add_textbox(Inches(0.95), Inches(5.45), Inches(10.8), Inches(0.7))
        orb_frame = orb_box.text_frame
        orb_run = orb_frame.paragraphs[0].add_run()
        orb_run.text = "Evidence orbs: " + (", ".join(slide_plan.evidence_orbs) if slide_plan.evidence_orbs else "reasoned synthesis")
        orb_run.font.size = Pt(12)
        orb_run.font.color.rgb = RGBColor(51, 65, 85)

        cue_panel = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(9.45),
            Inches(1.9),
            Inches(2.7),
            Inches(3.05),
        )
        cue_panel.fill.solid()
        cue_panel.fill.fore_color.rgb = RGBColor(*accent)
        cue_panel.fill.transparency = 0.16
        cue_panel.line.fill.background()
        cue_frame = cue_panel.text_frame
        cue_frame.word_wrap = True
        cue_title = cue_frame.paragraphs[0].add_run()
        cue_title.text = "Delivery cues"
        cue_title.font.size = Pt(13)
        cue_title.font.bold = True
        cue_title.font.color.rgb = RGBColor(255, 255, 255)
        for cue in (
            "Open with the stake, not the detail.",
            "Land one proof point before expanding.",
            "Pause before the final takeaway.",
        ):
            cue_paragraph = cue_frame.add_paragraph()
            cue_paragraph.text = cue
            cue_paragraph.font.size = Pt(11.5)
            cue_paragraph.font.color.rgb = RGBColor(255, 255, 255)
            cue_paragraph.space_after = Pt(8)

    def _paint_background(
        self,
        *,
        slide,
        accent: tuple[int, int, int],
        slide_width,
        slide_height,
        background_path: Path | None,
        soften_overlay: float,
    ) -> None:
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches

        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = RGBColor(12, 18, 28)

        if background_path:
            slide.shapes.add_picture(str(background_path), 0, 0, width=slide_width, height=slide_height)

        overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_width, slide_height)
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = RGBColor(248, 250, 252)
        overlay.fill.transparency = soften_overlay
        overlay.line.fill.background()

        top_band = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(0),
            slide_width,
            Inches(0.38),
        )
        top_band.fill.solid()
        top_band.fill.fore_color.rgb = RGBColor(*accent)
        top_band.fill.transparency = 0.08
        top_band.line.fill.background()

        orb = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(9.35),
            Inches(0.55),
            Inches(3.3),
            Inches(3.3),
        )
        orb.fill.solid()
        orb.fill.fore_color.rgb = RGBColor(*accent)
        orb.fill.transparency = 0.84
        orb.line.fill.background()

    def _render_visual_brief(self, *, slide, slide_plan, accent: tuple[int, int, int]) -> None:
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches, Pt

        visual_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(7.75),
            Inches(1.2),
            Inches(4.95),
            Inches(4.95),
        )
        visual_box.fill.solid()
        visual_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        visual_box.fill.transparency = 0.08
        visual_box.line.color.rgb = RGBColor(*accent)
        visual_frame = visual_box.text_frame
        visual_frame.word_wrap = True
        visual_header = visual_frame.paragraphs[0].add_run()
        visual_header.text = f"{slide_plan.visual_type.replace('_', ' ').title()} concept"
        visual_header.font.size = Pt(16)
        visual_header.font.bold = True
        visual_header.font.color.rgb = RGBColor(*accent)
        visual_body = visual_frame.add_paragraph()
        visual_body.text = slide_plan.visual_brief
        visual_body.font.size = Pt(14)
        visual_body.font.color.rgb = RGBColor(51, 65, 85)

    def _place_visual_asset(
        self,
        *,
        slide,
        slide_plan,
        context_map: dict[str, RetrievalResult],
        accent: tuple[int, int, int],
        used_visual_urls: set[str],
    ) -> bool:
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches, Pt

        image_urls = self._select_image_urls(
            slide_plan.evidence_orbs,
            context_map,
            limit=2,
            excluded_urls=used_visual_urls,
        )
        if not image_urls:
            image_urls = self._select_image_urls(
                slide_plan.evidence_orbs,
                context_map,
                limit=1,
                excluded_urls=set(),
            )
        if not image_urls:
            return False

        image_blobs: list[bytes] = []
        for image_url in image_urls:
            image_bytes = self._download_image(image_url)
            if image_bytes:
                image_blobs.append(image_bytes)

        if not image_blobs:
            return False

        used_visual_urls.update(image_urls[: len(image_blobs)])

        if len(image_blobs) == 1:
            slide.shapes.add_picture(BytesIO(image_blobs[0]), Inches(7.75), Inches(1.2), width=Inches(4.95), height=Inches(4.95))
        else:
            slide.shapes.add_picture(BytesIO(image_blobs[0]), Inches(7.75), Inches(1.2), width=Inches(4.95), height=Inches(2.38))
            slide.shapes.add_picture(BytesIO(image_blobs[1]), Inches(7.75), Inches(3.78), width=Inches(4.95), height=Inches(2.37))

        caption = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(8.05),
            Inches(5.72),
            Inches(4.35),
            Inches(0.42),
        )
        caption.fill.solid()
        caption.fill.fore_color.rgb = RGBColor(255, 255, 255)
        caption.fill.transparency = 0.12
        caption.line.color.rgb = RGBColor(*accent)
        caption_frame = caption.text_frame
        caption_frame.word_wrap = True
        caption_run = caption_frame.paragraphs[0].add_run()
        caption_run.text = slide_plan.visual_brief[:120]
        caption_run.font.size = Pt(10)
        caption_run.font.color.rgb = RGBColor(31, 41, 55)
        return True

    def _select_image_urls(
        self,
        evidence_orbs: list[str],
        context_map: dict[str, RetrievalResult],
        *,
        limit: int,
        excluded_urls: set[str],
    ) -> list[str]:
        selected: list[str] = []
        seen: set[str] = set()

        for orb in evidence_orbs:
            item = context_map.get(orb)
            if not item:
                continue
            for candidate in self._candidate_image_urls(item):
                if candidate not in seen and candidate not in excluded_urls:
                    selected.append(candidate)
                    seen.add(candidate)
                if len(selected) >= limit:
                    return selected

        for item in context_map.values():
            for candidate in self._candidate_image_urls(item):
                if candidate not in seen and candidate not in excluded_urls:
                    selected.append(candidate)
                    seen.add(candidate)
                if len(selected) >= limit:
                    return selected
        return selected

    def _candidate_image_urls(self, item: RetrievalResult) -> list[str]:
        candidates: list[str] = []
        metadata_url = item.metadata.get("image_url")
        if isinstance(metadata_url, str) and metadata_url.startswith(("http://", "https://")):
            candidates.append(metadata_url)
        content_url = self._extract_first_image_url(item.content)
        if content_url:
            candidates.append(content_url)
        return candidates

    @staticmethod
    def _download_image(url: str) -> bytes | None:
        try:
            response = httpx.get(url, timeout=20.0, follow_redirects=True)
            response.raise_for_status()
        except Exception:
            return None

        content_type = response.headers.get("content-type", "").lower()
        if not content_type.startswith("image/"):
            return None
        if "webp" in content_type or "svg" in content_type:
            return None
        return response.content

    def _select_background_path(self, topic: str) -> Path | None:
        if not self.backgrounds_dir or not self.backgrounds_dir.exists():
            return None
        preferred_background = self.backgrounds_dir / "2a1de542-04c7-45ed-bd9b-f68586de0987.jpeg"
        if preferred_background.exists():
            return preferred_background
        candidates = sorted(
            path
            for path in self.backgrounds_dir.iterdir()
            if path.is_file() and path.suffix.lower() in {".jpg", ".jpeg", ".png"}
        )
        if not candidates:
            return None
        index = int(sha1(topic.encode("utf-8")).hexdigest(), 16) % len(candidates)
        return candidates[index]

    @staticmethod
    def _extract_first_image_url(content: str) -> str | None:
        matches = re.findall(r"!\[[^\]]*\]\(([^)]+)\)", content)
        for match in matches:
            cleaned = match.strip().strip("<>").strip()
            if cleaned.startswith(("http://", "https://")) and "base64-image-removed" not in cleaned.lower():
                return cleaned
        return None

    @staticmethod
    def _orb_id(item: RetrievalResult) -> str:
        return f"ORB-{item.id[-8:]}"

    @staticmethod
    def _slugify(value: str) -> str:
        cleaned = re.sub(r"[^a-zA-Z0-9]+", "-", value).strip("-").lower()
        return cleaned[:80] or "presentation"
